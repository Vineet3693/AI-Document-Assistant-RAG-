"""
DocuMind AI - Document Loader
Load and extract text from various document formats (PDF, DOCX, TXT, CSV, XLSX, PPTX).
"""

import os
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass
from datetime import datetime

from ..utils.logger import get_module_logger
from ..utils.helpers import get_file_extension

logger = get_module_logger("document_loader")


@dataclass
class DocumentData:
    """Represents loaded document data."""
    filename: str
    file_path: str
    text: str
    pages: int
    metadata: Dict[str, Any]
    chunks: List[str]
    load_time_ms: int
    success: bool
    error_message: Optional[str] = None


class DocumentLoader:
    """
    Load and extract text from various document formats.
    
    Supported formats:
    - PDF (.pdf) - using PyPDF2
    - Word (.docx) - using python-docx
    - Text (.txt) - plain text
    - CSV (.csv) - using pandas
    - Excel (.xlsx) - using pandas/openpyxl
    - PowerPoint (.pptx) - using python-pptx
    - Markdown (.md) - plain text
    """
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.txt': self._load_txt,
            '.csv': self._load_csv,
            '.xlsx': self._load_xlsx,
            '.pptx': self._load_pptx,
            '.md': self._load_txt,
            '.json': self._load_txt,
            '.xml': self._load_txt,
        }
    
    def load(self, file_path: str) -> DocumentData:
        """
        Load a document from file path.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            DocumentData with extracted text and metadata
        """
        start_time = datetime.now()
        path = Path(file_path)
        
        if not path.exists():
            return DocumentData(
                filename=path.name,
                file_path=str(path),
                text="",
                pages=0,
                metadata={},
                chunks=[],
                load_time_ms=0,
                success=False,
                error_message=f"File not found: {file_path}",
            )
        
        ext = path.suffix.lower()
        
        if ext not in self.supported_extensions:
            return DocumentData(
                filename=path.name,
                file_path=str(path),
                text="",
                pages=0,
                metadata={},
                chunks=[],
                load_time_ms=0,
                success=False,
                error_message=f"Unsupported file type: {ext}",
            )
        
        try:
            logger.info(f"Loading document: {path.name} ({ext})")
            
            # Call appropriate loader
            loader_method = self.supported_extensions[ext]
            result = loader_method(path)
            
            end_time = datetime.now()
            load_time_ms = int((end_time - start_time).total_seconds() * 1000)
            
            logger.info(f"Successfully loaded {path.name}: {len(result['text'])} chars, "
                       f"{result['pages']} pages, {load_time_ms}ms")
            
            return DocumentData(
                filename=path.name,
                file_path=str(path),
                text=result['text'],
                pages=result['pages'],
                metadata=result.get('metadata', {}),
                chunks=[],  # Will be populated by TextSplitter
                load_time_ms=load_time_ms,
                success=True,
            )
            
        except Exception as e:
            logger.error(f"Error loading {path.name}: {str(e)}")
            end_time = datetime.now()
            load_time_ms = int((end_time - start_time).total_seconds() * 1000)
            
            return DocumentData(
                filename=path.name,
                file_path=str(path),
                text="",
                pages=0,
                metadata={},
                chunks=[],
                load_time_ms=load_time_ms,
                success=False,
                error_message=str(e),
            )
    
    def _load_pdf(self, path: Path) -> Dict:
        """Load PDF file using PyPDF2."""
        try:
            import PyPDF2
            
            text_pages = []
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_pages.append(page_text.strip())
            
            full_text = "\n\n".join(text_pages)
            
            # Extract metadata
            metadata = {}
            if reader.metadata:
                metadata = {
                    'title': reader.metadata.get('/Title', ''),
                    'author': reader.metadata.get('/Author', ''),
                    'subject': reader.metadata.get('/Subject', ''),
                    'creator': reader.metadata.get('/Creator', ''),
                    'producer': reader.metadata.get('/Producer', ''),
                    'creation_date': str(reader.metadata.get('/CreationDate', '')),
                }
            
            return {
                'text': full_text,
                'pages': num_pages,
                'metadata': metadata,
            }
            
        except Exception as e:
            logger.error(f"PDF loading error for {path.name}: {str(e)}")
            raise
    
    def _load_docx(self, path: Path) -> Dict:
        """Load Word document using python-docx."""
        try:
            from docx import Document
            
            doc = Document(path)
            paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            full_text = "\n\n".join(paragraphs)
            
            # Count approximate pages (assuming ~500 words per page)
            word_count = len(full_text.split())
            num_pages = max(1, word_count // 500)
            
            metadata = {
                'author': doc.core_properties.author if doc.core_properties.author else '',
                'title': doc.core_properties.title if doc.core_properties.title else '',
                'subject': doc.core_properties.subject if doc.core_properties.subject else '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
            }
            
            return {
                'text': full_text,
                'pages': num_pages,
                'metadata': metadata,
            }
            
        except Exception as e:
            logger.error(f"DOCX loading error for {path.name}: {str(e)}")
            raise
    
    def _load_txt(self, path: Path) -> Dict:
        """Load plain text file."""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                full_text = f.read()
            
            # Count approximate pages (assuming ~300 words per page for text files)
            word_count = len(full_text.split())
            num_pages = max(1, word_count // 300)
            
            return {
                'text': full_text,
                'pages': num_pages,
                'metadata': {},
            }
            
        except UnicodeDecodeError:
            # Try with latin-1 encoding as fallback
            with open(path, 'r', encoding='latin-1') as f:
                full_text = f.read()
            
            word_count = len(full_text.split())
            num_pages = max(1, word_count // 300)
            
            return {
                'text': full_text,
                'pages': num_pages,
                'metadata': {},
            }
    
    def _load_csv(self, path: Path) -> Dict:
        """Load CSV file using pandas."""
        try:
            import pandas as pd
            
            df = pd.read_csv(path)
            
            # Convert to readable text format
            text_parts = []
            text_parts.append(f"CSV File: {path.name}")
            text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
            text_parts.append(f"Rows: {len(df)}")
            text_parts.append("\nData:")
            
            # Convert dataframe to string representation
            text_parts.append(df.to_string())
            
            full_text = "\n".join(text_parts)
            
            return {
                'text': full_text,
                'pages': max(1, len(df) // 50),  # Approximate pages
                'metadata': {
                    'columns': df.columns.tolist(),
                    'row_count': len(df),
                },
            }
            
        except Exception as e:
            logger.error(f"CSV loading error for {path.name}: {str(e)}")
            raise
    
    def _load_xlsx(self, path: Path) -> Dict:
        """Load Excel file using pandas."""
        try:
            import pandas as pd
            
            # Read all sheets
            xls = pd.ExcelFile(path)
            text_parts = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                text_parts.append(f"\n{'='*50}")
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                text_parts.append(f"Rows: {len(df)}")
                text_parts.append("\nData:")
                text_parts.append(df.to_string())
            
            full_text = "\n".join(text_parts)
            total_rows = sum(pd.read_excel(path, sheet_name=sheet).shape[0] 
                           for sheet in xls.sheet_names)
            
            return {
                'text': full_text,
                'pages': max(1, total_rows // 50),
                'metadata': {
                    'sheets': xls.sheet_names,
                    'total_rows': total_rows,
                },
            }
            
        except Exception as e:
            logger.error(f"XLSX loading error for {path.name}: {str(e)}")
            raise
    
    def _load_pptx(self, path: Path) -> Dict:
        """Load PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation
            
            prs = Presentation(path)
            slide_texts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = [f"\n--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                
                slide_texts.append("\n".join(slide_text_parts))
            
            full_text = "\n\n".join(slide_texts)
            num_slides = len(prs.slides)
            
            return {
                'text': full_text,
                'pages': num_slides,
                'metadata': {
                    'slide_count': num_slides,
                },
            }
            
        except Exception as e:
            logger.error(f"PPTX loading error for {path.name}: {str(e)}")
            raise
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.supported_extensions.keys())
